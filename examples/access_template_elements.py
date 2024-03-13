from pathlib import Path
from typing import Final

import pptx

from pptx.presentation import Presentation


PROJECT_ROOT_PATH: Final[Path] = Path(__file__).parents[1]


def access_template_elements() -> Presentation:
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])

    title = slide.shapes.title
    title.text = "Hello, World!"

    text_frame = slide.shapes.placeholders[1]
    text_frame.text = "This is a text box."
    return presentation


if __name__ == "__main__":
    presentation = access_template_elements()
    presentation.save(PROJECT_ROOT_PATH / "dst.pptx")
