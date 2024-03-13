from pathlib import Path
from typing import Final

import pptx

from pptx.presentation import Presentation
from pptx.util import Inches


PROJECT_ROOT_PATH: Final[Path] = Path(__file__).parents[1]


def add_table() -> Presentation:
    table_data = [
        ["Name", "Age"],
        ["Alice", "20"],
        ["Bob", "21"],
        ["Charlie", "22"],
    ]

    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    row_size, column_size = len(table_data), len(table_data[0])
    table = slide.shapes.add_table(
        row_size,  # row size
        column_size,  # column size
        Inches(1),  # left
        Inches(1),  # top
        Inches(6),  # width
        Inches(3),  # height
    ).table
    for i in range(row_size):
        for j in range(column_size):
            table.cell(i, j).text = table_data[i][j]
    # カラムの長さは指定できる
    table.columns[0].width = Inches(4)
    table.columns[1].width = Inches(2)

    return presentation


if __name__ == "__main__":
    presentation = add_table()
    presentation.save(PROJECT_ROOT_PATH / "dst.pptx")
