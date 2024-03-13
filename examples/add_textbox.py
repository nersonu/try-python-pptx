from pathlib import Path
from typing import Final

import pptx

from pptx.dml.color import RGBColor
from pptx.enum.text import PP_PARAGRAPH_ALIGNMENT
from pptx.presentation import Presentation
from pptx.util import Inches, Pt


PROJECT_ROOT_PATH: Final[Path] = Path(__file__).parents[1]


def add_text_box() -> Presentation:
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    textbox = slide.shapes.add_textbox(
        Inches(1),  # left
        Inches(1),  # top
        Inches(6),  # width
        Inches(0.5),  # height
    )
    textbox.text = "This is a text box."
    for paragraph in textbox.text_frame.paragraphs:
        paragraph.alignment = PP_PARAGRAPH_ALIGNMENT.CENTER  # type: ignore
        for run in paragraph.runs:
            run.font.bold = True
            run.font.size = Pt(20)
            run.font.color.rgb = RGBColor(255, 0, 0)
    return presentation


if __name__ == "__main__":
    presentation = add_text_box()
    presentation.save(PROJECT_ROOT_PATH / "dst.pptx")
