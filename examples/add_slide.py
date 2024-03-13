from pathlib import Path
from typing import Final

import pptx

from pptx.presentation import Presentation


PROJECT_ROOT_PATH: Final[Path] = Path(__file__).parents[1]


def add_slide(template_path: Path) -> Presentation:
    presentation = pptx.Presentation(template_path)
    presentation.slides.add_slide(presentation.slide_layouts[0])
    return presentation


if __name__ == "__main__":
    template_path = PROJECT_ROOT_PATH / "data" / "template_example.pptx"
    presentation = add_slide(template_path)
    presentation.save(PROJECT_ROOT_PATH / "dst.pptx")
