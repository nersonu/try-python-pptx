from io import BytesIO
from pathlib import Path
from typing import Final

import matplotlib.pyplot as plt
import pptx

from pptx.presentation import Presentation
from pptx.util import Inches


PROJECT_ROOT_PATH: Final[Path] = Path(__file__).parents[1]


def add_figure(figure_path: Path) -> Presentation:
    presentation = pptx.Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.add_picture(
        str(figure_path),
        Inches(1),  # left
        Inches(1),  # top
    )
    plt.plot([*range(20)])
    with BytesIO() as stream:
        plt.savefig(stream, format="png")
        slide.shapes.add_picture(
            stream,
            Inches(3),  # left
            Inches(3),  # top
            width=Inches(3),  # width
        )
    return presentation


if __name__ == "__main__":
    figure_path = PROJECT_ROOT_PATH / "data" / "figure_example.png"
    presentation = add_figure(figure_path)
    presentation.save(PROJECT_ROOT_PATH / "dst.pptx")
